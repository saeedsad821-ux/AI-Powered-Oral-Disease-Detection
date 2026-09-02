"""Deliverables consistency scan: key numbers must match across artifacts."""

import json
import sys

from pptx import Presentation

sys.stdout.reconfigure(encoding="utf-8")

KEY = {
    "12,320 images": lambda t: ("12,320" in t or "12320" in t),
    "train 8,624": lambda t: ("8,624" in t or "8624" in t),
    "val/test 1,848": lambda t: ("1,848" in t or "1848" in t),
    "seed 42": lambda t: "seed 42" in t,
    "batch 32": lambda t: "32" in t,
    "224x224": lambda t: "224" in t,
    "weights 0.73..1.64": lambda t: ("0.73" in t and "1.64" in t),
    "6 classes": lambda t: "6" in t,
    "[0,255]": lambda t: ("[0,255]" in t or "0-255" in t or "0-255" in t),
}

docs = {}
for f in ["deliverables/report_Preprocessing.md",
          "deliverables/viva_preparation_Preprocessing.md",
          "deliverables/presentation_EDA.md",
          "deliverables/report_EDA.md", "README.md"]:
    docs[f] = open(f, encoding="utf-8").read()

p = Presentation("presentation/AI_Powered_Oral_Disease_Detection_System.pptx")
ppt = "\n".join(sh.text_frame.text for s in p.slides
                for sh in s.shapes if sh.has_text_frame)
for s in p.slides:
    if s.has_notes_slide:
        ppt += "\n" + s.notes_slide.notes_text_frame.text
docs["pptx"] = ppt

nb = json.load(open("notebooks/02_Preprocessing.ipynb", encoding="utf-8"))
nbt = ""
for c in nb["cells"]:
    s = c["source"]
    nbt += (s if isinstance(s, str) else "".join(s)) + "\n"
    for o in c.get("outputs", []):
        if o.get("output_type") == "stream":
            t = o.get("text", "")
            nbt += (t if isinstance(t, str) else "".join(t)) + "\n"
docs["notebook"] = nbt

print(f"{'artifact':<42}" + "".join(f"{k[:13]:>14}" for k in KEY))
ok = True
for name, text in docs.items():
    row = []
    for fn in KEY.values():
        v = fn(text)
        ok = ok and v
        row.append("Y" if v else "N")
    print(f"{name:<42}" + "".join(f"{v:>14}" for v in row))
print("\nALL CONSISTENT:", ok)
